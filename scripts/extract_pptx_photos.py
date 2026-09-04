#!/usr/bin/env python3
# Install dependencies first:
#   pip install python-pptx Pillow
"""
Extract every slide photo from a PowerPoint deck, in slide order, into a
web-ready gallery: resized/compressed JPEGs plus a JSON manifest that
gallery.js reads on videos_podcasts.html.

Exact-duplicate images (same file re-used on multiple slides) are written
once. Tiny images (icons/logos rather than photos) are skipped based on a
minimum pixel-dimension threshold.

Usage:
    python scripts/extract_pptx_photos.py Presentation2.pptx
    python scripts/extract_pptx_photos.py Presentation2.pptx --out-dir images/gallery --manifest data/gallery.json
"""

import argparse
import hashlib
import io
import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

REPO_ROOT = Path(__file__).resolve().parent.parent
MAX_DIMENSION = 1600
JPEG_QUALITY = 82
MIN_WIDTH = 300
MIN_HEIGHT = 300


def log(msg):
    print(f"[extract_pptx_photos] {msg}", file=sys.stderr)


def iter_slide_pictures(prs):
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield slide_index, shape


def resize_for_web(img):
    if img.mode not in ("RGB", "L"):
        img = img.convert("RGB")
    w, h = img.size
    if max(w, h) > MAX_DIMENSION:
        scale = MAX_DIMENSION / max(w, h)
        img = img.resize((round(w * scale), round(h * scale)), Image.LANCZOS)
    return img


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", help="Path to the source .pptx file")
    parser.add_argument("--out-dir", default=str(REPO_ROOT / "images" / "gallery"))
    parser.add_argument("--manifest", default=str(REPO_ROOT / "data" / "gallery.json"))
    parser.add_argument("--prefix", default="photo")
    args = parser.parse_args()

    out_dir = Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(args.pptx)

    seen_hashes = set()
    manifest = []
    skipped_small = 0
    skipped_dupe = 0
    skipped_error = 0

    for slide_index, shape in iter_slide_pictures(prs):
        try:
            # .blob avoids python-pptx's own format-detection (which chokes
            # on formats like MPO); we hand raw bytes straight to Pillow.
            raw = shape.image.blob
        except Exception as err:
            log(f"slide {slide_index}: could not read image ({err}); skipping")
            skipped_error += 1
            continue

        digest = hashlib.sha1(raw).hexdigest()
        if digest in seen_hashes:
            skipped_dupe += 1
            continue

        try:
            img = Image.open(io.BytesIO(raw))
            img.load()
        except Exception as err:
            log(f"slide {slide_index}: unreadable image data ({err}); skipping")
            skipped_error += 1
            continue

        w, h = img.size
        if w < MIN_WIDTH or h < MIN_HEIGHT:
            skipped_small += 1
            continue

        seen_hashes.add(digest)
        n = len(manifest) + 1
        filename = f"{args.prefix}-{n:03d}.jpg"
        web_img = resize_for_web(img)
        web_img.save(out_dir / filename, "JPEG", quality=JPEG_QUALITY, optimize=True)

        manifest.append({
            "file": f"images/gallery/{filename}",
            "slide": slide_index,
            "width": web_img.width,
            "height": web_img.height,
        })

    manifest_path = Path(args.manifest)
    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    manifest_path.write_text(
        json.dumps({"count": len(manifest), "photos": manifest}, indent=2),
        encoding="utf-8",
    )

    log(
        f"Wrote {len(manifest)} photos to {out_dir} "
        f"(skipped {skipped_dupe} duplicates, {skipped_small} too-small, {skipped_error} unreadable)"
    )
    log(f"Manifest: {manifest_path}")


if __name__ == "__main__":
    main()
